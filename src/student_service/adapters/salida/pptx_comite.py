"""Adapter de salida: informe para comité curricular (PPTX + PDF).

Usa el .pptx original de la coordinación como plantilla (conserva logos, franja
vinotinto y layout) y repuebla las tablas con los datos del repositorio:
  - Lámina 3: distribución de asistencia + lista "No asisten".
  - Lámina 4: matriz Nombre · A · NP · D · No asiste (dos columnas).
  - Lámina 5: incidencias (estudiantes que no superan la materia), si hay datos.

La plantilla contiene datos reales (PII) y NO va al repositorio: se configura
la ruta local con --plantilla o config.
"""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ...domain.asistencia import ResumenEstudiante, distribucion_por_presentes, no_asistentes

VINO = RGBColor(0x8A, 0x15, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
FONT = "Arial"
SOFFICE_CANDS = ["soffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"]


def _set_cell(cell, text, *, size=8, bold=False, color=DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    f = run.font
    f.size = Pt(size)
    f.name = FONT
    f.bold = bold
    f.color.rgb = color


def _del(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _find(slide, name):
    return next(s for s in slide.shapes if s.name == name)


def _style_matrix(tbl, headers, data, widths):
    tbl.first_row = False
    tbl.horz_banding = False
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        _set_cell(
            tbl.cell(0, j), h, size=8, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT, fill=VINO,
        )
    for i, row in enumerate(data, start=1):
        band = RGBColor(0xF2, 0xF4, 0xF6) if i % 2 == 0 else WHITE
        _set_cell(tbl.cell(i, 0), row[0], size=7.5, align=PP_ALIGN.LEFT, fill=band)
        for j, v in enumerate(row[1:], start=1):
            _set_cell(tbl.cell(i, j), v, size=7.5, align=PP_ALIGN.CENTER, fill=band)


class PptxInformeComite:
    def __init__(self, plantilla: str, exportar_pdf: bool = True) -> None:
        self.plantilla = plantilla
        self.exportar_pdf = exportar_pdf

    def render(
        self,
        resumenes: list[ResumenEstudiante],
        n_sesiones: int,
        salida: str,
        incidencias: dict[str, Any] | None = None,
    ) -> str:
        n = len(resumenes)
        dist = distribucion_por_presentes(resumenes, n_sesiones)
        no_asisten = no_asistentes(resumenes)
        prs = Presentation(self.plantilla)

        # rellenar placeholders de la plantilla limpia (docente / fechas)
        meta = (incidencias or {}).get("_plantilla", {})
        reemplazos = {
            "[Nombre del docente]": meta.get("docente", ""),
            "[dd/mm/aaaa – dd/mm/aaaa]": meta.get("fechas", ""),
        }
        for slide in prs.slides:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        for ph, val in reemplazos.items():
                            if val and ph in r.text:
                                r.text = r.text.replace(ph, val)

        def pct(c: int) -> str:
            return f"{round(100.0 * c / n, 2):.2f}".replace(".", ",") if n else "0,00"

        # ---- Lámina 3: tabla de distribución + lista "No asisten" ----
        s = prs.slides[2]
        dist_tbl = next(sh for sh in s.shapes if sh.has_table).table
        for ridx, k in enumerate(range(n_sesiones, -1, -1), start=1):
            etiqueta = "No asiste" if k == 0 else f"Asiste a {k} sesiones"
            for run in dist_tbl.cell(ridx, 0).text_frame.paragraphs[0].runs:
                run.text = ""
            _set_cell(dist_tbl.cell(ridx, 0), etiqueta, size=11, color=DARK)
            _set_cell(dist_tbl.cell(ridx, 1), dist[k], size=11, align=PP_ALIGN.CENTER)
            _set_cell(dist_tbl.cell(ridx, 2), pct(dist[k]), size=11, bold=True,
                      align=PP_ALIGN.CENTER)
        # la plantilla trae filas para 6 sesiones; si la cohorte tiene menos,
        # eliminar las sobrantes (quedarían con los datos del original)
        for tr in dist_tbl._tbl.tr_lst[n_sesiones + 2:]:
            dist_tbl._tbl.remove(tr)

        img = _find(s, "Imagen 2")
        left, top, w, h = img.left, img.top, img.width, img.height
        _del(img)
        gt = s.shapes.add_table(len(no_asisten) + 1, 1, left, top, w, h).table
        gt.first_row = False
        gt.horz_banding = False
        gt.columns[0].width = w
        _set_cell(gt.cell(0, 0), "Nombre", size=8, bold=True, color=WHITE, fill=VINO)
        for i, nom in enumerate(no_asisten, start=1):
            band = RGBColor(0xE9, 0xF1, 0xE6) if i % 2 else WHITE
            _set_cell(gt.cell(i, 0), nom, size=7.5, fill=band)

        # ---- Lámina 4: matriz A/NP/D en dos tablas ----
        s = prs.slides[3]
        headers = ["Nombre", "A", "NP", "D", "No asiste o tiempo incompleto"]
        widths = [2.55, 0.30, 0.35, 0.30, 0.94]
        mid = (len(resumenes) + 1) // 2
        for name, part in zip(("Imagen 4", "Imagen 7"), (resumenes[:mid], resumenes[mid:])):
            img = _find(s, name)
            left, top, w, h = img.left, img.top, img.width, img.height
            _del(img)
            scale = Emu(w).inches / sum(widths)
            tbl = s.shapes.add_table(len(part) + 1, 5, left, top, w, h).table
            data = [
                [r.nombre, r.presentes, r.ausentes, r.parciales, r.no_asiste_o_incompleto]
                for r in part
            ]
            _style_matrix(tbl, headers, data, [x * scale for x in widths])

        # ---- Lámina 5: incidencias (no superan la materia) ----
        reprobados = (incidencias or {}).get("no_superan")
        sld_ids = prs.slides._sldIdLst
        if reprobados:
            s = prs.slides[4]
            im_n, im_t = _find(s, "Imagen 4"), _find(s, "Imagen 7")
            left = min(im_n.left, im_t.left)
            top = im_n.top
            w = (im_t.left + im_t.width) - left
            _del(im_n)
            _del(im_t)
            tbl = s.shapes.add_table(len(reprobados) + 1, 2, left, top, w, Inches(0.35)).table
            tbl.first_row = False
            tbl.horz_banding = False
            tbl.columns[0].width = Inches(Emu(w).inches - 1.15)
            tbl.columns[1].width = Inches(1.15)
            _set_cell(tbl.cell(0, 0), "Nombre del estudiante", size=10, bold=True,
                      color=WHITE, fill=VINO)
            _set_cell(tbl.cell(0, 1), "Total", size=10, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER, fill=VINO)
            for i, r in enumerate(reprobados, start=1):
                band = RGBColor(0xF6, 0xEC, 0xEE) if i % 2 else WHITE
                _set_cell(tbl.cell(i, 0), r["nombre"], size=10, fill=band)
                _set_cell(tbl.cell(i, 1), f"{r['score']:.2f}%".replace(".", ","),
                          size=10, bold=True, align=PP_ALIGN.CENTER, fill=band)
            for sid in list(sld_ids)[6:4:-1]:  # comentarios (texto manual): fuera
                sld_ids.remove(sid)
        else:
            for sid in list(sld_ids)[6:3:-1]:  # sin calificaciones: fuera 5-7
                sld_ids.remove(sid)

        prs.save(salida)
        if self.exportar_pdf:
            self._a_pdf(salida)
        return salida

    @staticmethod
    def _a_pdf(pptx_path: str) -> str | None:
        soffice = next(
            (c for c in SOFFICE_CANDS if shutil.which(c) or Path(c).exists()), None
        )
        if not soffice:
            return None
        out = Path(pptx_path)
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out.parent), str(out)],
            check=True, capture_output=True,
        )
        return str(out.with_suffix(".pdf"))
